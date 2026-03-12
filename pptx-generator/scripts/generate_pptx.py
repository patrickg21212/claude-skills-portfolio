#!/usr/bin/env python3
"""
Generate branded .pptx presentations from structured text input.
Uses a brand config JSON to enforce fonts, colors, and layouts.

Usage:
    python generate_pptx.py --config brand_config.json --input content.md --output presentation.pptx
    python generate_pptx.py --config brand_config.json --input content.md --output presentation.pptx --template template.pptx
"""

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def parse_markdown_slides(text: str) -> list[dict]:
    """Parse markdown-style input into slide data structures.

    Expected format:
        # Slide Title
        Subtitle or body text

        ## Section Title
        - Bullet point one
        - Bullet point two
        - Bullet point three

        Body paragraph text goes here.

    Returns list of dicts with keys: title, subtitle, bullets, body, layout_hint
    """
    slides = []
    current_slide = None

    for line in text.strip().split("\n"):
        stripped = line.strip()

        # H1 = title slide
        if stripped.startswith("# ") and not stripped.startswith("## "):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "title": stripped[2:].strip(),
                "subtitle": "",
                "bullets": [],
                "body": "",
                "layout_hint": "title",
            }

        # H2 = content slide
        elif stripped.startswith("## "):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "title": stripped[3:].strip(),
                "subtitle": "",
                "bullets": [],
                "body": "",
                "layout_hint": "content",
            }

        # Bullet points
        elif stripped.startswith("- ") or stripped.startswith("* "):
            if current_slide is None:
                current_slide = {
                    "title": "",
                    "subtitle": "",
                    "bullets": [],
                    "body": "",
                    "layout_hint": "content",
                }
            current_slide["bullets"].append(stripped[2:].strip())

        # Subtitle (first non-bullet line after title slide)
        elif (
            current_slide
            and current_slide["layout_hint"] == "title"
            and not current_slide["subtitle"]
            and stripped
            and not current_slide["bullets"]
        ):
            current_slide["subtitle"] = stripped

        # Body text
        elif stripped and current_slide:
            if current_slide["body"]:
                current_slide["body"] += "\n" + stripped
            else:
                current_slide["body"] = stripped

    if current_slide:
        slides.append(current_slide)

    return slides


def apply_font(run, config: dict, style: str = "body"):
    """Apply brand font settings to a text run."""
    fonts = config.get("fonts", {})
    font_config = fonts.get(style, fonts.get("body", {}))

    run.font.name = font_config.get("name", "Calibri")
    run.font.size = Pt(font_config.get("size", 18))

    if font_config.get("bold"):
        run.font.bold = True
    if font_config.get("italic"):
        run.font.italic = True

    color = font_config.get("color")
    if color:
        run.font.color.rgb = hex_to_rgb(color)


def set_slide_background(slide, color_hex: str):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def create_title_slide(prs: Presentation, slide_data: dict, config: dict):
    """Create a title slide with brand styling."""
    colors = config.get("colors", {})
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)

    # Set background
    bg_color = colors.get("primary_bg", "#FFFFFF")
    set_slide_background(slide, bg_color)

    # Title
    if slide.placeholders[0] is not None:
        title_shape = slide.placeholders[0]
        title_shape.text = slide_data["title"]
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                apply_font(run, config, "heading")
            paragraph.alignment = PP_ALIGN.LEFT

    # Subtitle
    if len(slide.placeholders) > 1 and slide_data.get("subtitle"):
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = slide_data["subtitle"]
        for paragraph in subtitle_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                apply_font(run, config, "subtitle")
            paragraph.alignment = PP_ALIGN.LEFT

    return slide


def create_content_slide(prs: Presentation, slide_data: dict, config: dict):
    """Create a content slide with title, bullets, and/or body text."""
    colors = config.get("colors", {})
    layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(layout)

    # Background
    bg_color = colors.get("slide_bg", "#FFFFFF")
    set_slide_background(slide, bg_color)

    # Title
    if slide.placeholders[0] is not None and slide_data.get("title"):
        title_shape = slide.placeholders[0]
        title_shape.text = slide_data["title"]
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                apply_font(run, config, "heading")

    # Content area
    if len(slide.placeholders) > 1:
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.clear()

        # Add bullets
        if slide_data.get("bullets"):
            for i, bullet in enumerate(slide_data["bullets"]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = Pt(6)
                for run in p.runs:
                    apply_font(run, config, "body")

                # Accent color for bullet character
                accent = colors.get("accent")
                if accent:
                    p.font.color.rgb = hex_to_rgb(
                        colors.get("body_text", "#333333")
                    )

        # Add body text
        elif slide_data.get("body"):
            p = tf.paragraphs[0]
            p.text = slide_data["body"]
            for run in p.runs:
                apply_font(run, config, "body")

    return slide


def generate_presentation(
    content_path: str,
    config_path: str,
    output_path: str,
    template_path: str | None = None,
):
    """Main function: parse input, apply brand, generate .pptx."""

    # Load brand config
    with open(config_path) as f:
        config = json.load(f)

    # Load content
    with open(content_path) as f:
        content = f.read()

    # Parse slides from markdown
    slides_data = parse_markdown_slides(content)

    if not slides_data:
        print("Error: No slides found in input. Use # for title slides, ## for content slides.", file=sys.stderr)
        sys.exit(1)

    # Create presentation (from template or blank)
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # Set slide dimensions (widescreen 16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Generate each slide
    for slide_data in slides_data:
        if slide_data["layout_hint"] == "title":
            create_title_slide(prs, slide_data, config)
        else:
            create_content_slide(prs, slide_data, config)

    # Save
    prs.save(output_path)
    print(f"Generated {len(slides_data)} slides -> {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate branded .pptx from markdown")
    parser.add_argument("--config", required=True, help="Path to brand config JSON")
    parser.add_argument("--input", required=True, help="Path to markdown content file")
    parser.add_argument("--output", required=True, help="Output .pptx file path")
    parser.add_argument("--template", default=None, help="Optional .pptx template to use as base")
    args = parser.parse_args()

    generate_presentation(args.input, args.config, args.output, args.template)
